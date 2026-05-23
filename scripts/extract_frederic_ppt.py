#!/usr/bin/env python3
"""Extrae texto del PPT del curso SEM (Frederic Marimon) a docs/."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

DEFAULT_PPT = Path(
    "/Users/claudiolarrea/Library/CloudStorage/OneDrive-Personal/"
    "11 Investigacion/Ecuaciones estructurales/"
    "Structural Equation Modeling (SEM).pptx"
)
OUT = Path(__file__).resolve().parents[1] / "docs" / "frederic_sem_ppt_extract.txt"


def main(ppt_path: Path = DEFAULT_PPT) -> None:
    prs = Presentation(str(ppt_path))
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            s.text.strip()
            for s in slide.shapes
            if hasattr(s, "text") and s.text.strip()
        ]
        if texts:
            blocks.append(f"\n=== SLIDE {i} ===\n" + "\n".join(texts))
    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text("\n".join(blocks), encoding="utf-8")
    print(f"Wrote {len(prs.slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
